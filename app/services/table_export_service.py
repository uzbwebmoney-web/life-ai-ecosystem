from __future__ import annotations

import csv
import io
import re
from typing import Literal

TableFormat = Literal["csv", "xlsx", "md", "html"]
PresentationFormat = Literal["pptx", "docx"]


def extract_table_from_text(text: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in (text or "").splitlines():
        line = line.strip()
        if not line:
            continue
        if "|" in line and line.count("|") >= 2:
            cells = [c.strip() for c in line.strip("|").split("|")]
            if cells and not all(set(c) <= {"-", " "} for c in cells):
                rows.append(cells)
            continue
        if "\t" in line:
            rows.append([c.strip() for c in line.split("\t")])
            continue
        m = re.match(r"^[\d\-•*]+\.?\s*(.+?):\s*(.+)$", line)
        if m:
            rows.append([m.group(1).strip(), m.group(2).strip()])
    return rows


def build_table_bytes(
    rows: list[list[str]],
    *,
    fmt: TableFormat = "csv",
    title: str = "table",
) -> tuple[bytes, str, str]:
    safe = re.sub(r'[<>:"/\\|?*]', "_", title).strip() or "table"
    if fmt == "csv":
        buf = io.StringIO()
        writer = csv.writer(buf)
        for row in rows:
            writer.writerow(row)
        data = buf.getvalue().encode("utf-8-sig")
        return data, f"{safe}.csv", "text/csv"

    if fmt == "md":
        lines = []
        if rows:
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
            for row in rows[1:]:
                lines.append("| " + " | ".join(row) + " |")
        content = "\n".join(lines).encode("utf-8")
        return content, f"{safe}.md", "text/markdown"

    if fmt == "html":
        body = ["<table border='1' cellpadding='4'>"]
        for i, row in enumerate(rows):
            tag = "th" if i == 0 else "td"
            body.append("<tr>" + "".join(f"<{tag}>{cell}</{tag}>" for cell in row) + "</tr>")
        body.append("</table>")
        html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{safe}</title></head><body>{''.join(body)}</body></html>"
        return html.encode("utf-8"), f"{safe}.html", "text/html"

    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = safe[:31]
        for row in rows:
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue(), f"{safe}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    except ImportError:
        return build_table_bytes(rows, fmt="csv", title=title)


def _split_slide_bullets(body: str) -> list[str]:
    lines: list[str] = []
    for raw in body.splitlines():
        line = raw.strip()
        if not line:
            continue
        line = re.sub(r"^[\d\-•*]+\.?\s*", "", line)
        lines.append(line[:500])
    return lines[:12]


def build_pptx_bytes(title: str, slides: list[tuple[str, str]]) -> tuple[bytes, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title[:200]
    try:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slides[0][1][:300] if slides else ""
    except Exception:
        pass

    content_layout = prs.slide_layouts[1]
    for slide_title, body in slides:
        if not slide_title.strip():
            continue
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide_title[:200]
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        bullets = _split_slide_bullets(body)
        if not bullets:
            bullets = [body[:500]]
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    safe = re.sub(r'[<>:"/\\|?*]', "_", title).strip() or "presentation"
    return buf.getvalue(), f"{safe}.pptx"


def build_presentation_bytes(
    title: str,
    slides: list[tuple[str, str]],
    *,
    fmt: PresentationFormat = "pptx",
) -> tuple[bytes, str]:
    if fmt == "pptx":
        return build_pptx_bytes(title, slides)
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for slide_title, body in slides:
        doc.add_heading(slide_title, level=1)
        for line in body.splitlines():
            if line.strip():
                doc.add_paragraph(line.strip())
    buf = io.BytesIO()
    doc.save(buf)
    safe = re.sub(r'[<>:"/\\|?*]', "_", title).strip() or "presentation"
    return buf.getvalue(), f"{safe}.docx"


def slides_from_text(text: str, *, default_title: str = "Slide") -> list[tuple[str, str]]:
    slides: list[tuple[str, str]] = []
    current_title = default_title
    current_lines: list[str] = []
    for line in (text or "").splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("## "):
            if current_lines:
                slides.append((current_title, "\n".join(current_lines)))
            current_title = stripped[3:].strip()[:200]
            current_lines = []
        elif stripped.startswith("# "):
            current_title = stripped[2:].strip()[:200]
        else:
            current_lines.append(stripped)
    if current_lines:
        slides.append((current_title, "\n".join(current_lines)))
    if not slides and text.strip():
        slides.append((default_title, text.strip()))
    return slides
